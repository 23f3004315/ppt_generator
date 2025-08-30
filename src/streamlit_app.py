import streamlit as st
import google.generativeai as genai
import openai
import anthropic
import markdown
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import io
import json
import re
import time
from typing import Dict, List, Optional, Any

# --- App Configuration ---
st.set_page_config(
    page_title="Your Text, Your Style – AI Presentation Generator",
    page_icon="✨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Configuration constants
MAX_FILE_SIZE_MB = 50
MAX_TEXT_LENGTH = 100000
MAX_SLIDES = 20
API_RETRY_COUNT = 3
API_TIMEOUT = 30

# Custom CSS for better UI
st.markdown("""
<style>
    .success-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #d4edda;
        border: 1px solid #c3e6cb;
        color: #155724;
        margin: 1rem 0;
    }
    
    .warning-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #fff3cd;
        border: 1px solid #ffeaa7;
        color: #856404;
        margin: 1rem 0;
    }
    
    .info-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #d1ecf1;
        border: 1px solid #bee5eb;
        color: #0c5460;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# --- Helper Functions ---

def validate_file_upload(uploaded_file) -> bool:
    """Validate uploaded PowerPoint file."""
    if uploaded_file is None:
        return False
    
    # Check file size
    if uploaded_file.size > MAX_FILE_SIZE_MB * 1024 * 1024:
        st.error(f"❌ File too large! Maximum size is {MAX_FILE_SIZE_MB}MB. Your file is {uploaded_file.size / (1024*1024):.1f}MB.")
        return False
    
    # Check file extension
    allowed_extensions = ['.pptx', '.potx']
    if not any(uploaded_file.name.lower().endswith(ext) for ext in allowed_extensions):
        st.error("❌ Invalid file type! Please upload a .pptx or .potx file.")
        return False
    
    return True

def validate_text_input(text: str) -> bool:
    """Validate text input length and content."""
    if not text or not text.strip():
        return False
    
    if len(text) > MAX_TEXT_LENGTH:
        st.error(f"❌ Text too long! Maximum length is {MAX_TEXT_LENGTH:,} characters. Your text is {len(text):,} characters.")
        return False
    
    return True

def api_call_with_retry(api_function, *args, **kwargs):
    """Execute API call with retry logic."""
    for attempt in range(API_RETRY_COUNT):
        try:
            return api_function(*args, **kwargs)
        except Exception as e:
            if attempt == API_RETRY_COUNT - 1:
                raise e
            else:
                st.warning(f"API attempt {attempt + 1} failed, retrying... ({str(e)[:100]})")
                time.sleep(2 ** attempt)  # Exponential backoff
    return None

def clean_layout_name(name: str) -> str:
    """Clean and standardize layout names for better matching."""
    return re.sub(r'^Layout \d+:\s*[\'"]?|[\'"]?$', '', name).strip()

def extract_template_images(template_bytes: bytes) -> Dict[str, bytes]:
    """Extract images from the template for reuse."""
    images = {}
    try:
        template_stream = io.BytesIO(template_bytes)
        prs = Presentation(template_stream)
        
        # Extract images from slides
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_key = f"slide_{slide_idx}_shape_{shape_idx}"
                        images[image_key] = image_bytes
                    except:
                        continue
        
        # Extract images from slide layouts
        for layout_idx, layout in enumerate(prs.slide_layouts):
            for shape_idx, shape in enumerate(layout.shapes):
                if hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_key = f"layout_{layout_idx}_shape_{shape_idx}"
                        images[image_key] = image_bytes
                    except:
                        continue
                        
    except Exception as e:
        st.warning(f"Could not extract images from template: {str(e)}")
    
    return images

def place_template_images(slide, template_images: Dict[str, bytes], slide_index: int):
    """Place appropriate template images back into slides."""
    try:
        # Look for image placeholders in the slide
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format'):
                # This is a placeholder that might accept images
                placeholder_type = shape.placeholder_format.type
                # Type 18 is typically picture placeholder
                if placeholder_type == 18:  # Picture placeholder
                    # Try to find a suitable image from template
                    image_key = f"slide_{slide_index}_shape_0"
                    if image_key in template_images:
                        try:
                            image_stream = io.BytesIO(template_images[image_key])
                            shape.insert_picture(image_stream)
                        except:
                            # If insertion fails, continue without image
                            pass
    except Exception:
        # If image placement fails, continue without images
        pass

def preserve_template_formatting(slide, original_layout):
    """Preserve original template formatting as much as possible."""
    try:
        # Attempt to preserve background and color scheme
        if hasattr(slide, 'background') and hasattr(original_layout, 'background'):
            try:
                slide.background = original_layout.background
            except:
                pass
        
        # Preserve any master slide elements
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # Try to maintain original font properties
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.runs:
                        # Preserve template font settings where possible
                        try:
                            # The template's font settings are usually preserved automatically
                            # by using the template's slide layouts
                            pass
                        except:
                            continue
    except Exception:
        # If formatting preservation fails, continue with default
        pass

def get_ppt_layouts(template_bytes: bytes) -> Dict[str, Any]:
    """
    Analyzes PowerPoint file bytes and returns a dictionary of its layouts.
    """
    if not template_bytes:
        return {}
    
    try:
        template_stream = io.BytesIO(template_bytes)
        prs = Presentation(template_stream)
        
        layouts = {}
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = f"Layout {i}: '{layout.name}'"
            clean_name = clean_layout_name(layout_name)
            
            # Skip blank layouts and add layout info
            if "blank" not in clean_name.lower():
                layouts[layout_name] = {
                    'layout_obj': layout,
                    'clean_name': clean_name,
                    'placeholders': [p.name for p in layout.placeholders if hasattr(p, 'name')]
                }
        
        return layouts
        
    except Exception as e:
        st.error(f"Error reading template file: {str(e)}")
        return {}

def process_markdown_content(content: str) -> str:
    """Convert markdown to plain text while preserving structure."""
    try:
        # Convert markdown to HTML first
        html = markdown.markdown(content)
        # Remove HTML tags for plain text
        clean_text = re.sub(r'<[^>]+>', '', html)
        # Clean up extra whitespace
        clean_text = re.sub(r'\n\s*\n', '\n\n', clean_text)
        return clean_text.strip()
    except:
        # Fallback to original content if markdown processing fails
        return content

def generate_with_gemini(api_key: str, user_text: str, guidance: str, layout_names: List[str]) -> Optional[List[Dict]]:
    """
    Calls the Google Gemini API to structure the text into presentation slides.
    """
    def _make_gemini_call():
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        # Clean layout names for the prompt
        clean_layouts = [clean_layout_name(name) for name in layout_names]
        
        system_prompt = f"""
        You are an expert presentation creator. Your task is to analyze the provided text and structure it into a series of presentation slides.
        
        Available slide layouts: {', '.join(clean_layouts)}
        
        IMPORTANT: Structure your response as a valid JSON object with a "slides" key containing an array of slide objects.
        
        Each slide object must have:
        1. "layout": Choose from the available layouts above
        2. "content": Object with "title" (string) and "body" (array of bullet points)
        
        CRITICAL CONTENT GUIDELINES:
        - Create 4-12 slides depending on content length and complexity
        - Keep titles under 60 characters (2 lines max)
        - Limit body content to 4-6 bullet points per slide
        - Each bullet point must be a simple string. DO NOT use markdown, nested dictionaries, or complex JSON objects for a single bullet point.
        - Each bullet point should be 1-2 sentences maximum (under 100 characters)
        - Break long content into multiple slides rather than cramming everything
        - Use clear, concise language suitable for slide presentation
        - First slide should be a title/cover slide with just title and subtitle
        - Avoid dense paragraphs - use short, punchy bullet points
        - If content is extensive, create more slides with less content each
        """
        
        user_prompt = f"Text to convert into presentation:\n\n{user_text}\n\n"
        if guidance:
            user_prompt += f"Additional guidance: {guidance}"
        
        response = model.generate_content(
            [system_prompt, user_prompt],
            generation_config=genai.types.GenerationConfig(
                response_mime_type="application/json",
                temperature=0.7
            )
        )
        
        return response
    
    try:
        response = api_call_with_retry(_make_gemini_call)
        
        try:
            response_data = json.loads(response.text)
            slides = response_data.get('slides', [])
            
            if not slides:
                st.error("AI generated empty slides. Please try again with different text.")
                return None
                
            return slides
            
        except json.JSONDecodeError as e:
            st.error(f"AI response was not valid JSON. Please try again. Error: {str(e)}")
            st.code(response.text)
            return None
            
    except Exception as e:
        st.error(f"Gemini API error: {str(e)}")
        return None

def generate_with_anthropic(api_key: str, user_text: str, guidance: str, layout_names: List[str]) -> Optional[List[Dict]]:
    """
    Calls the Anthropic Claude API to structure the text into presentation slides.
    """
    try:
        client = anthropic.Anthropic(api_key=api_key)
        
        # Clean layout names for the prompt
        clean_layouts = [clean_layout_name(name) for name in layout_names]
        
        system_prompt = f"""
        You are an expert presentation creator. Analyze the provided text and structure it into presentation slides.
        
        Available layouts: {', '.join(clean_layouts)}
        
        Return a JSON object with "slides" key containing an array of slide objects.
        Each slide needs:
        1. "layout": Pick from available layouts
        2. "content": Object with "title" (string) and "body" (array of strings)
        
        CRITICAL CONTENT GUIDELINES:
        - Create 4-12 slides depending on content length and complexity
        - Keep titles under 60 characters (2 lines max)
        - Limit body content to 4-6 bullet points per slide
        - Each bullet point should be 1-2 sentences maximum (under 100 characters)
        - Each bullet point must be a simple string. DO NOT use markdown, nested dictionaries, or complex JSON objects for a single bullet point.
        - Break long content into multiple slides rather than cramming everything
        - Use clear, concise language suitable for slide presentation
        - First slide should be a title/cover slide
        - Avoid dense paragraphs - use short, punchy bullet points
        - If content is extensive, create more slides with less content each
        """
        
        user_prompt = f"Convert this text into a presentation:\n\n{user_text}"
        if guidance:
            user_prompt += f"\n\nGuidance: {guidance}"
        
        response = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=4000,
            temperature=0.7,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}]
        )
        
        try:
            # Extract JSON from Claude's response
            response_text = response.content[0].text
            
            # Try to find JSON in the response
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            if json_match:
                json_str = json_match.group()
            else:
                json_str = response_text
            
            response_data = json.loads(json_str)
            slides = response_data.get('slides', [])
            
            if not slides:
                st.error("AI generated empty slides. Please try again.")
                return None
                
            return slides
            
        except json.JSONDecodeError as e:
            st.error(f"Claude response was not valid JSON: {str(e)}")
            st.code(response_text)
            return None
            
    except anthropic.AuthenticationError:
        st.error("❌ Authentication Error: Invalid Anthropic API key.")
        return None
    except anthropic.RateLimitError:
        st.error("❌ Rate limit exceeded. Please try again later.")
        return None
    except Exception as e:
        st.error(f"Anthropic API error: {str(e)}")
        return None

def generate_with_openai(api_key: str, user_text: str, guidance: str, layout_names: List[str]) -> Optional[List[Dict]]:
    """
    Calls the OpenAI API to structure the text into presentation slides.
    """
    try:
        client = openai.OpenAI(api_key=api_key)
        
        # Clean layout names for the prompt
        clean_layouts = [clean_layout_name(name) for name in layout_names]
        
        system_prompt = f"""
        You are an expert presentation creator. Analyze the provided text and structure it into presentation slides.
        
        Available layouts: {', '.join(clean_layouts)}
        
        Return a JSON object with "slides" key containing an array of slide objects.
        Each slide needs:
        1. "layout": Pick from available layouts
        2. "content": Object with "title" (string) and "body" (array of strings)
        
        CRITICAL CONTENT GUIDELINES:
        - Create 4-12 slides depending on content length and complexity
        - Keep titles under 60 characters (2 lines max)
        - Limit body content to 4-6 bullet points per slide
        - Each bullet point should be 1-2 sentences maximum (under 100 characters)
        - Each bullet point must be a simple string. DO NOT use markdown, nested dictionaries, or complex JSON objects for a single bullet point.
        - Break long content into multiple slides rather than cramming everything
        - Use clear, concise language suitable for slide presentation
        - First slide should be a title/cover slide
        - Avoid dense paragraphs - use short, punchy bullet points
        - If content is extensive, create more slides with less content each
        """
        
        user_prompt = f"Convert this text into a presentation:\n\n{user_text}"
        if guidance:
            user_prompt += f"\n\nGuidance: {guidance}"
        
        response = client.chat.completions.create(
            model="gpt-4o-nano",
            response_format={"type": "json_object"},
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.7
        )
        
        try:
            response_data = json.loads(response.choices[0].message.content)
            slides = response_data.get('slides', [])
            
            if not slides:
                st.error("AI generated empty slides. Please try again.")
                return None
                
            return slides
            
        except json.JSONDecodeError as e:
            st.error(f"AI response was not valid JSON: {str(e)}")
            return None
            
    except openai.AuthenticationError:
        st.error("❌ Authentication Error: Invalid OpenAI API key.")
        return None
    except openai.RateLimitError:
        st.error("❌ Rate limit exceeded. Please try again later.")
        return None
    except Exception as e:
        st.error(f"OpenAI API error: {str(e)}")
        return None

def find_best_layout_match(requested_layout: str, available_layouts: Dict) -> Optional[Any]:
    """Find the best matching layout from available layouts."""
    
    # First, try exact match with cleaned names
    for layout_key, layout_info in available_layouts.items():
        if layout_info['clean_name'].lower() == requested_layout.lower():
            return layout_info['layout_obj']
    
    # Try partial matches
    requested_lower = requested_layout.lower()
    
    # Common layout type mappings
    layout_mappings = {
        'title': ['title', 'cover', 'intro'],
        'content': ['content', 'bullet', 'text'],
        'two': ['two', 'comparison', 'split'],
        'section': ['section', 'divider', 'header']
    }
    
    for layout_key, layout_info in available_layouts.items():
        layout_name_lower = layout_info['clean_name'].lower()
        
        # Check for keyword matches
        for keyword, synonyms in layout_mappings.items():
            if keyword in requested_lower:
                for synonym in synonyms:
                    if synonym in layout_name_lower:
                        return layout_info['layout_obj']
    
    # Fallback: return first available layout
    if available_layouts:
        return list(available_layouts.values())[0]['layout_obj']
    
    return None

def validate_and_optimize_slides(slides_data: List[Dict]) -> List[Dict]:
    """
    Validates and optimizes slide content to prevent overflow and improve readability.
    """
    optimized_slides = []
    
    for i, slide_info in enumerate(slides_data):
        content = slide_info.get('content', {})
        title = content.get('title', '')
        body = content.get('body', [])
        
        # Optimize title
        if len(title) > 60:
            title = title[:57] + "..."
        
        # Optimize body content
        if isinstance(body, list):
            optimized_body = []
            for point in body[:6]:  # Max 6 points per slide
                point_text = str(point).strip()
                # Split very long points
                if len(point_text) > 100:
                    # Try to split at sentence boundary
                    sentences = point_text.split('. ')
                    if len(sentences) > 1:
                        optimized_body.append(sentences[0] + '.')
                        if len(optimized_body) < 6:
                            remaining = '. '.join(sentences[1:])
                            if len(remaining) > 100:
                                remaining = remaining[:97] + "..."
                            optimized_body.append(remaining)
                    else:
                        optimized_body.append(point_text[:97] + "...")
                else:
                    optimized_body.append(point_text)
            
            # If we have too much content, suggest splitting
            if len(body) > 6:
                # Split into multiple slides
                chunks = [body[i:i+5] for i in range(0, len(body), 5)]
                base_title = title.replace(" (Part 1)", "").replace(" - Part 1", "")
                
                for idx, chunk in enumerate(chunks):
                    slide_copy = slide_info.copy()
                    if idx == 0:
                        slide_copy['content'] = {'title': title, 'body': chunk}
                    else:
                        slide_copy['content'] = {'title': f"{base_title} (Part {idx+1})", 'body': chunk}
                    optimized_slides.append(slide_copy)
                continue
        else:
            optimized_body = [str(body)[:100] + "..." if len(str(body)) > 100 else str(body)]
        
        # Update slide content
        optimized_slide = slide_info.copy()
        optimized_slide['content'] = {'title': title, 'body': optimized_body}
        optimized_slides.append(optimized_slide)
    
    return optimized_slides

def generate_speaker_notes(api_key: str, llm_provider: str, slide_content: Dict, guidance: str = "") -> str:
    """Generate speaker notes for a slide using LLM."""
    try:
        title = slide_content.get('title', '')
        body = slide_content.get('body', [])
        body_text = '\n'.join(body) if isinstance(body, list) else str(body)
        
        prompt = f"""Generate concise speaker notes for this presentation slide:

Title: {title}
Content: {body_text}

{f"Presentation context: {guidance}" if guidance else ""}

Create 2-3 sentences of speaker notes that expand on the slide content, provide context, or suggest talking points. Keep it professional and helpful for the presenter."""

        if llm_provider == "OpenAI":
            client = openai.OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=200,
                temperature=0.7
            )
            return response.choices[0].message.content.strip()
            
        elif llm_provider == "Anthropic":
            client = anthropic.Anthropic(api_key=api_key)
            response = client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=200,
                temperature=0.7,
                messages=[{"role": "user", "content": prompt}]
            )
            return response.content[0].text.strip()
            
        elif llm_provider == "Gemini":
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel('gemini-1.5-flash-latest')
            response = model.generate_content(prompt)
            return response.text.strip()
            
    except Exception as e:
        return f"Speaker notes could not be generated: {str(e)}"

def create_presentation(slides_data: List[Dict], template_bytes: bytes, available_layouts: Dict, 
                       api_key: str = None, llm_provider: str = None, guidance: str = "", 
                       generate_notes: bool = False) -> Optional[io.BytesIO]:
    """
    Creates a PowerPoint presentation from structured slide data with enhanced template preservation.
    """
    try:
        template_stream = io.BytesIO(template_bytes)
        prs = Presentation(template_stream)
        
        # Extract template images for potential reuse
        template_images = extract_template_images(template_bytes)
        
        slides_created = 0
        
        for i, slide_info in enumerate(slides_data):
            try:
                layout_name = slide_info.get('layout', '')
                content = slide_info.get('content', {})
                
                # Find best matching layout
                slide_layout = find_best_layout_match(layout_name, available_layouts)
                
                if not slide_layout:
                    st.warning(f"No suitable layout found for slide {i+1}. Skipping.")
                    continue
                
                # Create slide
                slide = prs.slides.add_slide(slide_layout)
                slides_created += 1
                
                # Preserve template formatting
                preserve_template_formatting(slide, slide_layout)
                
                # Place template images where appropriate
                place_template_images(slide, template_images, i)
                
                # Fill placeholders with better content management
                title_set = False
                body_set = False
                
                for shape in slide.placeholders:
                    if not hasattr(shape, 'text_frame'):
                        continue
                        
                    shape_name = getattr(shape, 'name', '').lower()
                    
                    # Handle title placeholders
                    if not title_set and ('title' in shape_name) and content.get('title'):
                        title = content['title']
                        # Truncate title if too long
                        if len(title) > 60:
                            title = title[:57] + "..."
                        shape.text = title
                        title_set = True
                    
                    # Handle body/content placeholders with better content management
                    elif not body_set and any(keyword in shape_name for keyword in ['body', 'content', 'text']) and content.get('body'):
                        text_frame = shape.text_frame
                        text_frame.clear()
                        text_frame.word_wrap = True
                        
                        body_points = content.get('body', [])
                        if isinstance(body_points, list) and body_points:
                            # Limit the number of bullet points to prevent overflow
                            max_points = 6
                            limited_points = body_points[:max_points]
                            
                            if limited_points:
                                # Set first point
                                first_point = str(limited_points[0])
                                if len(first_point) > 100:
                                    first_point = first_point[:97] + "..."
                                
                                p = text_frame.paragraphs[0]
                                p.text = first_point
                                p.level = 0
                                
                                # Add remaining points as bullet points
                                for point in limited_points[1:]:
                                    point_text = str(point)
                                    if len(point_text) > 100:
                                        point_text = point_text[:97] + "..."
                                    
                                    p = text_frame.add_paragraph()
                                    p.text = point_text
                                    p.level = 0
                        
                        body_set = True
                
                # Generate speaker notes if requested
                if generate_notes and api_key and llm_provider:
                    try:
                        notes = generate_speaker_notes(api_key, llm_provider, content, guidance)
                        if notes and hasattr(slide, 'notes_slide'):
                            slide.notes_slide.notes_text_frame.text = notes
                    except Exception as e:
                        # Silently continue if notes generation fails
                        pass
                        
            except Exception as e:
                st.warning(f"Error creating slide {i+1}: {str(e)}")
                continue
        
        if slides_created == 0:
            st.error("No slides were successfully created.")
            return None
        
        # Save presentation to memory
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        return ppt_stream
        
    except Exception as e:
        st.error(f"Failed to create presentation: {str(e)}")
        return None

# --- Initialize Session State ---
if 'template_bytes' not in st.session_state:
    st.session_state.template_bytes = None
if 'layouts' not in st.session_state:
    st.session_state.layouts = {}
if 'template_name' not in st.session_state:
    st.session_state.template_name = None

# --- UI Layout ---

st.title("✨ Your Text, Your Style – AI Presentation Generator")
st.markdown("Transform any text into a beautifully formatted PowerPoint presentation using your own template and AI.")

# --- Sidebar for Configuration ---
with st.sidebar:
    st.header("⚙️ Configuration")
    
    # LLM Provider Selection
    llm_provider = st.selectbox(
        "Choose AI Provider", 
        ["OpenAI", "Gemini", "Anthropic"],
        help="Select your preferred AI service"
    )
    
    # API Key Input
    api_key = st.text_input(
        f"Enter your {llm_provider} API Key", 
        type="password",
        help=f"Your {llm_provider} API key (never stored)",
        key=f"api_key_{llm_provider}"
    )
    
    # Show API key status
    if api_key:
        st.success("✅ API Key provided")
    
    st.markdown("---")
    
    # Instructions
    st.header("📋 How It Works")
    st.markdown("""
    1. **Upload Template**: Choose a .pptx or .potx file
    2. **Add Content**: Paste your text content
    3. **Optional Guidance**: Specify tone or style
    4. **Generate**: AI creates structured slides
    5. **Download**: Get your formatted presentation
    """)
    
    st.markdown("---")
    st.info("🔒 Your API key is only used for this session and never stored.")

# --- Main Content Area ---

# Step 1: Template Upload
st.header("1️⃣ Upload PowerPoint Template")

uploaded_template = st.file_uploader(
    "Choose a PowerPoint template file (.pptx or .potx)",
    type=['pptx', 'potx'],
    help=f"Maximum file size: {MAX_FILE_SIZE_MB}MB. This template's layouts, fonts, and colors will be applied to your generated presentation",
    key='template_uploader'
)

# Validate and process template
if uploaded_template is not None:
    if validate_file_upload(uploaded_template):
        # Check if this is a new file
        if (st.session_state.template_name != uploaded_template.name or 
            st.session_state.template_bytes is None):
            
            with st.spinner(f"📋 Analyzing template: {uploaded_template.name}"):
                try:
                    # Store template data
                    st.session_state.template_bytes = uploaded_template.getvalue()
                    st.session_state.template_name = uploaded_template.name
                    
                    # Analyze layouts
                    st.session_state.layouts = get_ppt_layouts(st.session_state.template_bytes)
                    
                    if st.session_state.layouts:
                        st.success(f"✅ Template analyzed successfully! Found {len(st.session_state.layouts)} usable layouts.")
                    else:
                        st.error("❌ Could not find usable layouts in the template.")
                        
                except Exception as e:
                    st.error(f"❌ Error processing template: {str(e)}")
                    st.session_state.template_bytes = None
                    st.session_state.layouts = {}

# Display template info if available
if st.session_state.layouts:
    with st.expander("📋 Template Layouts Available", expanded=False):
        st.write("AI will choose from these layouts:")
        for i, (layout_name, layout_info) in enumerate(st.session_state.layouts.items(), 1):
            st.write(f"**{i}.** {layout_info['clean_name']}")

# Step 2: Content Input
st.header("2️⃣ Add Your Content")

col1, col2 = st.columns([2, 1])

with col1:
    input_text = st.text_area(
        "Paste your text content here (supports markdown)",
        height=300,
        max_chars=MAX_TEXT_LENGTH,
        placeholder=f"Paste your article, report, notes, or any text content here...\n\nSupports:\n• Plain text\n• Markdown formatting\n• Long-form prose\n\nMaximum length: {MAX_TEXT_LENGTH:,} characters\n\nThe AI will automatically break this into slides with appropriate titles and bullet points.",
        help="Supports plain text, markdown, or any written content",
        key="content_input"
    )
    
    # Process markdown if detected
    if input_text and ('**' in input_text or '#' in input_text or '*' in input_text):
        processed_text = process_markdown_content(input_text)
        if processed_text != input_text:
            st.info("📝 Markdown formatting detected and will be processed")
    
    # Character count
    if input_text:
        char_count = len(input_text)
        word_count = len(input_text.split())
        st.caption(f"📝 {char_count:,} characters, ~{word_count:,} words")
        
        # Content length guidance
        if word_count > 2000:
            st.warning("⚠️ **Large content detected!** Consider breaking this into multiple presentations or providing more specific guidance to help AI create focused slides.")
        elif word_count > 1000:
            st.info("💡 **Tip:** With substantial content, consider specifying the presentation type (e.g., 'executive summary', 'technical deep-dive') for better slide organization.")

with col2:
    st.subheader("3️⃣ Optional Guidance")
    guidance_text = st.text_input(
        "Presentation style/tone",
        placeholder="e.g., 'executive summary', 'technical deep-dive', 'sales pitch'",
        help="Optional: Guide the AI on tone, structure, or intended audience",
        key="guidance_input"
    )
    
    # Pre-defined guidance options
    st.write("**Quick Options:**")
    guidance_options = [
        "Executive summary",
        "Sales pitch", 
        "Investor pitch deck",
        "Technical presentation",
        "Training material",
        "Project update",
        "Research findings",
        "Visual-heavy presentation"
    ]
    
    # Create columns for better layout of buttons
    cols = st.columns(2)
    for i, option in enumerate(guidance_options):
        with cols[i % 2]:
            if st.button(option, key=f"guidance_btn_{i}", use_container_width=True):
                st.session_state.guidance_input = option
                st.rerun()
    
    st.markdown("---")
    
    # Optional speaker notes generation
    st.subheader("🗣️ Speaker Notes")
    generate_notes = st.checkbox(
        "Generate speaker notes for each slide",
        help="Uses AI to create helpful speaker notes for presentation delivery",
        key="generate_notes_checkbox"
    )

# Update guidance_text if it was set via button
if 'guidance_input' in st.session_state and st.session_state.guidance_input:
    guidance_text = st.session_state.guidance_input

# Step 4: Generation
st.header("4️⃣ Generate Presentation")

# Validation and generation
can_generate = all([
    api_key and len(api_key.strip()) > 0,
    st.session_state.template_bytes,
    input_text and len(input_text.strip()) > 0 and validate_text_input(input_text),
    st.session_state.layouts
])

if not can_generate:
    missing_items = []
    if not api_key or len(api_key.strip()) == 0:
        missing_items.append("API key")
    if not st.session_state.template_bytes:
        missing_items.append("PowerPoint template")
    if not input_text or len(input_text.strip()) == 0:
        missing_items.append("text content")
    if not st.session_state.layouts:
        missing_items.append("valid template layouts")
    
    st.warning(f"⚠️ Please provide: {', '.join(missing_items)}")
else:
    st.success("✅ Ready to generate! All requirements met.")

generate_col1, generate_col2, generate_col3 = st.columns([1, 2, 1])

with generate_col2:
    generate_button = st.button(
        "🚀 Generate Presentation",
        type="primary",
        use_container_width=True,
        disabled=not can_generate
    )

if generate_button and can_generate:
    # Progress tracking
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    try:
        # Step 1: AI Processing
        status_text.text("🤖 AI is analyzing your content...")
        progress_bar.progress(25)
        
        # Process markdown content if needed
        processed_content = input_text
        if '**' in input_text or '#' in input_text or '*' in input_text:
            processed_content = process_markdown_content(input_text)
        
        layout_names = list(st.session_state.layouts.keys())
        slides = None
        
        if llm_provider == "Gemini":
            slides = generate_with_gemini(api_key, processed_content, guidance_text, layout_names)
        elif llm_provider == "OpenAI":
            slides = generate_with_openai(api_key, processed_content, guidance_text, layout_names)
        elif llm_provider == "Anthropic":
            slides = generate_with_anthropic(api_key, processed_content, guidance_text, layout_names)
        
        if not slides:
            st.error("❌ Failed to generate slide structure. Please try again.")
            st.stop()
        
        # Optimize slides to prevent content overflow
        slides = validate_and_optimize_slides(slides)
        
        progress_bar.progress(50)
        status_text.text(f"✅ Generated {len(slides)} slides structure")
        
        # Display slide preview
        with st.expander("📋 Preview Generated Slides", expanded=True):
            st.info(f"✅ Generated {len(slides)} slides with optimized content to prevent overflow")
            for i, slide in enumerate(slides, 1):
                slide_title = slide.get('content', {}).get('title', 'Untitled')
                st.write(f"**Slide {i}: {slide_title}**")
                
                body = slide.get('content', {}).get('body', [])
                if body:
                    points_shown = 0
                    for point in body:
                        if points_shown >= 4:  # Show max 4 points in preview
                            st.write(f"• ... and {len(body) - points_shown} more points")
                            break
                        st.write(f"• {point}")
                        points_shown += 1
                    
                    # Show content optimization info
                    if len(body) > 6:
                        st.caption(f"⚠️ Content optimized: Originally {len(body)} points, showing top 6")
                
                if i < len(slides):
                    st.write("---")
        
        # Step 2: Create Presentation
        status_text.text("📊 Creating PowerPoint presentation...")
        progress_bar.progress(75)
        
        presentation_stream = create_presentation(
            slides, 
            st.session_state.template_bytes, 
            st.session_state.layouts,
            api_key=api_key,
            llm_provider=llm_provider,
            guidance=guidance_text,
            generate_notes=st.session_state.get('generate_notes_checkbox', False)
        )
        
        if not presentation_stream:
            st.error("❌ Failed to create presentation file.")
            st.stop()
        
        progress_bar.progress(100)
        status_text.text("✅ Presentation ready for download!")
        
        # Success and download
        st.success(f"🎉 Successfully created presentation with {len(slides)} slides!")
        
        # Download button
        st.download_button(
            label="📥 Download Presentation (.pptx)",
            data=presentation_stream,
            file_name=f"AI_Generated_Presentation_{st.session_state.template_name.split('.')[0]}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary"
        )
        
        st.balloons()
        
    except Exception as e:
        st.error(f"❌ An unexpected error occurred: {str(e)}")
    
    finally:
        progress_bar.empty()
        status_text.empty()

# --- Footer ---
st.markdown("---")
st.markdown(
    """
    <div style='text-align: center; color: #666; font-size: 0.8em;'>
        Made with ❤️ using Streamlit | Your API keys are never stored or logged
    </div>
    """, 
    unsafe_allow_html=True
)
